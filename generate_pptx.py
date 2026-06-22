import os
from pptx import Presentation
from pptx.util import Inches

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    assets_dir = '/Users/mac/Downloads/noorpath/assets/'
    
    # List of images in order
    images = []
    
    # Add Book Pages first
    for i in range(1, 8):
        img_path = os.path.join(assets_dir, f'book_page_{i}_welcome.png' if i==1 else f'book_page_{i}_alif.png' if i==2 else f'book_page_{i}_ba.png' if i==3 else f'book_page_{i}_ta.png' if i==4 else f'book_page_{i}_game.png' if i==5 else f'book_page_{i}_story.png' if i==6 else f'book_page_{i}_reward.png')
        # Fix: I used specific names in GenerateImage, let's just check if they exist
        potential_names = [
            f'book_page_{i}_welcome.png',
            f'book_page_{i}_alif.png',
            f'book_page_{i}_ba.png',
            f'book_page_{i}_ta.png',
            f'book_page_{i}_game.png',
            f'book_page_{i}_story.png',
            f'book_page_{i}_reward.png'
        ]
        for name in potential_names:
            p = os.path.join(assets_dir, name)
            if os.path.exists(p):
                images.append(p)
                break

    # Add Alphabet Cards
    alphabet_files = sorted([f for f in os.listdir(assets_dir) if f.startswith('alphabet_') and f.endswith('.png')], 
                            key=lambda x: int(x.split('_')[1]))
    for f in alphabet_files:
        images.append(os.path.join(assets_dir, f))

    for img_path in images:
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image to fill the slide
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_path = '/Users/mac/Downloads/noorpath/NoorPath_Academy_Lessons.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_presentation()

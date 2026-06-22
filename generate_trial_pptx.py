import os
from pptx import Presentation
from pptx.util import Inches

def create_trial_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    assets_dir = 'assets/trial_class/'
    
    # List of images in order
    images = [
        'welcome.png',
        'intro.png',
        'alif.png',
        'ba.png',
        'ta.png',
        'game.png',
        'moral.png',
        'reward.png'
    ]
    
    for img_name in images:
        img_path = os.path.join(assets_dir, img_name)
        if not os.path.exists(img_path):
            print(f"Warning: {img_path} not found. Skipping.")
            continue
            
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Add image to fill the slide
        slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    output_path = 'Trial_Class_3-4_Years.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_trial_presentation()

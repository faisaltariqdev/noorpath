import os
from pptx import Presentation
from pptx.util import Inches

def create_full_trial_presentation():
    prs = Presentation()
    
    # Set slide dimensions to 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Day 1
    day1_assets = 'assets/trial_class/'
    day1_images = [
        'welcome.png',
        'intro.png',
        'alif.png',
        'ba.png',
        'ta.png',
        'game.png',
        'moral.png',
        'reward.png'
    ]
    
    # Day 2
    day2_assets = 'assets/trial_class/day2/'
    day2_images = [
        'welcome_day2.png',
        'review_day1.png',
        'tha.png',
        'jeem.png',
        'ha.png',
        'game_day2.png',
        'moral_day2.png'
    ]
    
    # Day 3
    day3_assets = 'assets/trial_class/day3/'
    day3_images = [
        'welcome_day3.png',
        'kha.png',
        'dal.png',
        'dhal.png',
        'game_day3.png',
        'moral_day3.png',
        'graduation.png'
    ]
    
    def add_slides(assets_dir, images):
        for img_name in images:
            img_path = os.path.join(assets_dir, img_name)
            if not os.path.exists(img_path):
                print(f"Warning: {img_path} not found. Skipping.")
                continue
                
            slide_layout = prs.slide_layouts[6]  # Blank layout
            slide = prs.slides.add_slide(slide_layout)
            slide.shapes.add_picture(img_path, 0, 0, width=prs.slide_width, height=prs.slide_height)

    add_slides(day1_assets, day1_images)
    add_slides(day2_assets, day2_images)
    add_slides(day3_assets, day3_images)

    output_path = 'Full_3_Day_Trial_Class.pptx'
    prs.save(output_path)
    print(f"Presentation saved to {output_path}")

if __name__ == "__main__":
    create_full_trial_presentation()
